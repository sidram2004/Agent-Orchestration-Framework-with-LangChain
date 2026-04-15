import sqlite3
import functools
import time as time_module
from flask import Flask, render_template, request, redirect, url_for, session, jsonify, send_from_directory
from werkzeug.security import generate_password_hash, check_password_hash
from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation
from orchestration import run_workflow
from agents import TOOLS

print("[Startup] Initializing Flask app...")
app = Flask(__name__)
app.secret_key = "orchest_ai_secret_2024"

print("[Startup] App created. Registering routes...")


@app.template_global()
def get_badge_class(routing: str) -> str:
    r = (routing or "").upper()
    if "SIMPLE"   in r: return "badge-simple"
    if "TOOL"     in r: return "badge-tool"
    if "CONTENT"  in r: return "badge-content"
    if "MEDICAL"  in r: return "badge-medical"
    if "DECISION" in r: return "badge-decision"
    if "DEBUG"    in r: return "badge-debug"
    if "SHOPPING" in r: return "badge-shopping"
    return "badge-general"


def get_db():
    conn = sqlite3.connect("database.db")
    conn.row_factory = sqlite3.Row
    return conn


def login_required(f):
    @functools.wraps(f)
    def decorated(*args, **kwargs):
        if 'user' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated


def init_db():
    with get_db() as db:
        # Create tables if they don't exist
        db.execute("""CREATE TABLE IF NOT EXISTS users(
            id INTEGER PRIMARY KEY, username TEXT UNIQUE, password TEXT)""")
        db.execute("""CREATE TABLE IF NOT EXISTS chats(
            id INTEGER PRIMARY KEY, chat_id TEXT, username TEXT,
            message TEXT, response TEXT)""")
        db.execute("""CREATE TABLE IF NOT EXISTS user_memory(
            id INTEGER PRIMARY KEY, username TEXT, memory_key TEXT UNIQUE,
            memory_value TEXT, updated_at DATETIME DEFAULT CURRENT_TIMESTAMP)""")
        db.execute("""CREATE TABLE IF NOT EXISTS uploaded_files(
            id INTEGER PRIMARY KEY, username TEXT, chat_id TEXT,
            filename TEXT, filetype TEXT, uploaded_at DATETIME DEFAULT CURRENT_TIMESTAMP)""")

        # Migrate existing tables — safely add new columns if missing
        migrations = [
            ("users",  "language",      "TEXT DEFAULT 'en'"),
            ("users",  "created_at",    "DATETIME DEFAULT CURRENT_TIMESTAMP"),
            ("chats",  "routing",       "TEXT DEFAULT ''"),
            ("chats",  "response_time", "REAL DEFAULT 0"),
            ("chats",  "pinned",        "INTEGER DEFAULT 0"),
            ("chats",  "timestamp",     "DATETIME DEFAULT CURRENT_TIMESTAMP"),
        ]
        for table, col, col_type in migrations:
            try:
                db.execute(f"ALTER TABLE {table} ADD COLUMN {col} {col_type}")
            except Exception:
                pass  # Column already exists — skip
        db.commit()

init_db()


def get_user_memory_context(username):
    try:
        db = get_db()
        rows = db.execute("SELECT memory_key, memory_value FROM user_memory WHERE username=?", (username,)).fetchall()
        if not rows:
            return ""
        mem_str = "\n".join([f"- {r['memory_key']}: {r['memory_value']}" for r in rows])
        return f"\n### USER CORE CONTEXT (PERSISTENT MEMORY) ###\n{mem_str}\n"
    except Exception as e:
        print(f"[Memory Fetch Error] {e}")
        return ""


def extract_text_from_file(file):
    name = file.filename.lower()
    try:
        if name.endswith('.pdf'):
            reader = PdfReader(file)
            return "".join([p.extract_text() or "" for p in reader.pages]).strip()
        elif name.endswith('.docx'):
            doc = docx.Document(file)
            return "\n".join([p.text for p in doc.paragraphs]).strip()
        elif name.endswith(('.xlsx', '.xls')):
            return pd.read_excel(file).to_string(index=False)
        elif name.endswith('.csv'):
            return pd.read_csv(file).to_string(index=False)
        elif name.endswith('.pptx'):
            prs = Presentation(file)
            return "\n".join([shape.text for slide in prs.slides
                              for shape in slide.shapes if hasattr(shape, "text")]).strip()
        elif name.endswith('.txt'):
            return file.read().decode('utf-8').strip()
        return ""
    except Exception as e:
        return f"[Error extracting: {name}]"


# ── Auth ──────────────────────────────────────────────────────────────────────
@app.route('/')
def index():
    if 'user' in session: return redirect(url_for('chat_interface'))
    return render_template('landing.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        u, p = request.form['username'], request.form['password']
        try:
            db = get_db()
            db.execute("INSERT INTO users (username, password) VALUES (?,?)", (u, generate_password_hash(p)))
            db.commit()
            return redirect(url_for('login'))
        except sqlite3.IntegrityError:
            return "Username already exists. <a href='/register'>Try again</a>", 400
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        u, p = request.form.get('username'), request.form.get('password')
        db  = get_db()
        user = db.execute("SELECT * FROM users WHERE username=?", (u,)).fetchone()
        if user and check_password_hash(user['password'], p):
            session['user'] = u
            try:
                session['language'] = user['language'] or 'en'
            except (IndexError, KeyError):
                session['language'] = 'en'
            return redirect(url_for('chat_interface'))
        return "Invalid credentials. <a href='/login'>Try again</a>", 401
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))


# ── Chat ──────────────────────────────────────────────────────────────────────
@app.route('/chat')
@login_required
def chat_interface():
    current_chat = request.args.get('chat_id', 'Chat 1')
    db = get_db()
    rows = db.execute("""SELECT chat_id, MAX(id) as last_id FROM chats
        WHERE username=? GROUP BY chat_id ORDER BY last_id DESC LIMIT 20""", (session['user'],)).fetchall()
    chat_list = [r["chat_id"] for r in rows] or ["Chat 1"]
    
    # Ensure current workspace is at the top
    if current_chat in chat_list:
        chat_list.remove(current_chat)
    chat_list.insert(0, current_chat)
    
    # Strictly limit to top 6 for "minimum" history look
    chat_list = chat_list[:6]
    history   = db.execute(
        "SELECT id, message, response, routing, response_time, pinned FROM chats WHERE username=? AND chat_id=? ORDER BY id",
        (session['user'], current_chat)).fetchall()
    user = db.execute("SELECT language FROM users WHERE username=?", (session['user'],)).fetchone()
    try:
        lang = user['language'] if user else 'en'
        lang = lang or 'en'
    except (IndexError, KeyError):
        lang = 'en'
    return render_template('index.html', chat_list=chat_list, current_chat=current_chat,
                           history=history, username=session['user'], language=lang)


@app.route('/send_message', methods=['POST'])
@login_required
def send_message():
    chat_id    = request.form['chat_id']
    user_input = request.form.get('message', '').strip()
    db = get_db()
    if 'pdf' in request.files and request.files['pdf'].filename:
        file     = request.files['pdf']
        doc_text = extract_text_from_file(file)
        ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else 'unknown'
        db.execute("INSERT INTO uploaded_files (username, chat_id, filename, filetype) VALUES (?,?,?,?)",
                   (session['user'], chat_id, file.filename, ext))
        db.commit()
        if doc_text:
            user_input = f"{user_input}\n\n[Document: {file.filename}]\n{doc_text}"

    count = db.execute("SELECT COUNT(*) as c FROM chats WHERE username=? AND chat_id=?",
                       (session['user'], chat_id)).fetchone()['c']

    t_start = time_module.time()
    mem_ctx = get_user_memory_context(session['user'])
    result  = run_workflow(user_input, is_new_chat=(count == 0), core_memory=mem_ctx)
    elapsed = round(time_module.time() - t_start, 2)

    ai_resp    = result.get("answer", "")
    routing    = result.get("routing", "")
    tools_used = result.get("tools_used", [])
    suggested  = result.get("suggested_name")

    cursor = db.execute(
        "INSERT INTO chats (chat_id, username, message, response, routing, response_time) VALUES (?,?,?,?,?,?)",
        (chat_id, session['user'], user_input, ai_resp, routing, elapsed))
    db.commit()

    return jsonify({"id": cursor.lastrowid, "message": user_input, "answer": ai_resp,
                    "response": ai_resp, "routing": routing, "email": "",
                    "tools_used": tools_used, "suggested_name": suggested, "response_time": elapsed})


@app.route('/edit_message', methods=['POST'])
@login_required
def edit_message():
    msg_id  = request.form['msg_id']
    new_text = request.form['message']
    t_start = time_module.time()
    mem_ctx = get_user_memory_context(session['user'])
    result  = run_workflow(new_text, is_new_chat=False, core_memory=mem_ctx)
    elapsed = round(time_module.time() - t_start, 2)
    new_resp = result.get("answer", "")
    db = get_db()
    try:
        db.execute("UPDATE chats SET message=?, response=?, routing=?, response_time=? WHERE id=? AND username=?",
                   (new_text, new_resp, result.get("routing",""), elapsed, msg_id, session['user']))
        db.commit()
        return jsonify({"success": True, "new_response": new_resp,
                        "tools_used": result.get("tools_used",[]),
                        "routing": result.get("routing",""), "response_time": elapsed})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/rename_workspace', methods=['POST'])
@login_required
def rename_workspace():
    old_id, new_id = request.form['old_id'], request.form['new_id']
    db = get_db()
    try:
        db.execute("UPDATE chats SET chat_id=? WHERE username=? AND chat_id=?", (new_id, session['user'], old_id))
        db.commit()
        return jsonify({"success": True, "new_id": new_id})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/clear_chat', methods=['POST'])
@login_required
def clear_chat():
    chat_id = request.form['chat_id']
    db = get_db()
    try:
        db.execute("DELETE FROM chats WHERE username=? AND chat_id=?", (session['user'], chat_id))
        db.commit()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/delete_chat/<chat_id>')
@login_required
def delete_chat(chat_id):
    db = get_db()
    db.execute("DELETE FROM chats WHERE username=? AND chat_id=?", (session['user'], chat_id))
    db.commit()
    return redirect(url_for('chat_interface'))


# ── Feature 4: Search ─────────────────────────────────────────────────────────
@app.route('/search_history')
@login_required
def search_history():
    q = request.args.get('q', '').strip()
    if not q: return jsonify({"results": []})
    db = get_db()
    rows = db.execute("""SELECT id, chat_id, message, response, routing, timestamp
        FROM chats WHERE username=? AND (message LIKE ? OR response LIKE ?)
        ORDER BY id DESC LIMIT 20""", (session['user'], f'%{q}%', f'%{q}%')).fetchall()
    return jsonify({"results": [dict(r) for r in rows]})


# ── Feature 7: Pin Message ────────────────────────────────────────────────────
@app.route('/pin_message', methods=['POST'])
@login_required
def pin_message():
    msg_id = request.form['msg_id']
    state  = int(request.form.get('state', 1))
    db = get_db()
    try:
        db.execute("UPDATE chats SET pinned=? WHERE id=? AND username=?", (state, msg_id, session['user']))
        db.commit()
        return jsonify({"success": True, "pinned": state})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/pinned_messages')
@login_required
def pinned_messages():
    db = get_db()
    rows = db.execute(
        "SELECT id, chat_id, message, response, routing, timestamp FROM chats WHERE username=? AND pinned=1 ORDER BY id DESC",
        (session['user'],)).fetchall()
    return jsonify({"pinned": [dict(r) for r in rows]})


# ── Feature 6: User Memory ────────────────────────────────────────────────────
@app.route('/user_memory', methods=['GET', 'POST'])
@login_required
def user_memory():
    db = get_db()
    if request.method == 'GET':
        rows = db.execute(
            "SELECT memory_key, memory_value, updated_at FROM user_memory WHERE username=? ORDER BY updated_at DESC",
            (session['user'],)).fetchall()
        return jsonify({"memory": [dict(r) for r in rows]})
    key   = request.form.get('key','').strip()
    value = request.form.get('value','').strip()
    if key and value:
        try:
            db.execute("INSERT INTO user_memory (username,memory_key,memory_value) VALUES(?,?,?)",
                       (session['user'], key, value))
        except sqlite3.IntegrityError:
            db.execute("UPDATE user_memory SET memory_value=?,updated_at=CURRENT_TIMESTAMP WHERE username=? AND memory_key=?",
                       (value, session['user'], key))
        db.commit()
    return jsonify({"success": True})


@app.route('/delete_memory', methods=['POST'])
@login_required
def delete_memory():
    key = request.form.get('key')
    db  = get_db()
    db.execute("DELETE FROM user_memory WHERE username=? AND memory_key=?", (session['user'], key))
    db.commit()
    return jsonify({"success": True})


@app.route('/clear_all_memory', methods=['POST'])
@login_required
def clear_all_memory():
    db = get_db()
    db.execute("DELETE FROM user_memory WHERE username=?", (session['user'],))
    db.commit()
    return jsonify({"success": True})


# ── Feature 10: Profile ───────────────────────────────────────────────────────
@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    db = get_db()
    if request.method == 'POST':
        action = request.form.get('action')
        if action == 'change_password':
            old_p = request.form.get('old_password')
            new_p = request.form.get('new_password')
            user  = db.execute("SELECT * FROM users WHERE username=?", (session['user'],)).fetchone()
            if user and check_password_hash(user['password'], old_p):
                db.execute("UPDATE users SET password=? WHERE username=?",
                           (generate_password_hash(new_p), session['user']))
                db.commit()
                return jsonify({"success": True, "message": "Password updated!"})
            return jsonify({"success": False, "message": "Old password incorrect."}), 400
        elif action == 'change_language':
            lang = request.form.get('language', 'en')
            db.execute("UPDATE users SET language=? WHERE username=?", (lang, session['user']))
            db.commit()
            session['language'] = lang
            return jsonify({"success": True, "language": lang})

    stats = db.execute("""SELECT COUNT(*) as total_messages,
        COUNT(DISTINCT chat_id) as total_chats,
        AVG(response_time) as avg_response_time,
        MIN(timestamp) as member_since
        FROM chats WHERE username=?""", (session['user'],)).fetchone()
    files  = db.execute("SELECT COUNT(*) as c FROM uploaded_files WHERE username=?", (session['user'],)).fetchone()
    pinned = db.execute("SELECT COUNT(*) as c FROM chats WHERE username=? AND pinned=1", (session['user'],)).fetchone()
    user   = db.execute("SELECT language, created_at FROM users WHERE username=?", (session['user'],)).fetchone()
    return jsonify({
        "username": session['user'],
        "total_messages":    stats['total_messages'],
        "total_chats":       stats['total_chats'],
        "avg_response_time": round(stats['avg_response_time'] or 0, 2),
        "member_since":      stats['member_since'] or user['created_at'],
        "files_uploaded":    files['c'],
        "pinned_messages":   pinned['c'],
        "language":          user['language'] or 'en',
    })


# ── Analytics API ────────────────────────────────────────────────────────────
@app.route('/api/analytics')
@login_required
def analytics():
    db = get_db()
    total_users    = db.execute("SELECT COUNT(*) as c FROM users").fetchone()['c']
    total_chats    = db.execute("SELECT COUNT(DISTINCT chat_id) as c FROM chats").fetchone()['c']
    total_messages = db.execute("SELECT COUNT(*) as c FROM chats").fetchone()['c']
    avg_time       = db.execute("SELECT AVG(response_time) as a FROM chats").fetchone()['a']

    # Recent 10 workflows
    recent = db.execute(
        "SELECT username, chat_id, routing, response_time, timestamp FROM chats ORDER BY id DESC LIMIT 10"
    ).fetchall()

    # Pipeline distribution — normalize routing labels into clean categories
    all_rows = db.execute("SELECT routing FROM chats WHERE routing IS NOT NULL AND routing != ''").fetchall()
    dist_map = {}
    for row in all_rows:
        r = (row['routing'] or '').upper()
        if   'SIMPLE'   in r: key = 'Simple'
        elif 'TOOL'     in r: key = 'Tool'
        elif 'CONTENT'  in r: key = 'Content'
        elif 'MEDICAL'  in r: key = 'Medical'
        elif 'DECISION' in r: key = 'Decision'
        elif 'DEBUG'    in r: key = 'Debug'
        elif 'SHOPPING' in r: key = 'Shopping'
        elif 'GENERAL'  in r: key = 'General'
        elif 'ERROR'    in r: key = 'Error'
        else:                 key = 'General'
        dist_map[key] = dist_map.get(key, 0) + 1

    distribution = [{"routing": k, "cnt": v} for k, v in sorted(dist_map.items(), key=lambda x: -x[1])]

    # Daily trend last 7 days
    daily = db.execute("""
        SELECT DATE(timestamp) as day, COUNT(*) as cnt
        FROM chats WHERE timestamp >= DATE('now', '-7 days')
        GROUP BY DATE(timestamp) ORDER BY day ASC
    """).fetchall()

    # Top workspaces
    top_ws = db.execute("""
        SELECT chat_id, COUNT(*) as cnt FROM chats
        WHERE username=? GROUP BY chat_id ORDER BY cnt DESC LIMIT 6
    """, (session['user'],)).fetchall()

    return jsonify({
        "total_users":       total_users,
        "total_chats":       total_chats,
        "total_messages":    total_messages,
        "avg_response_time": round(avg_time or 0, 2),
        "recent":            [dict(r) for r in recent],
        "distribution":      distribution,
        "daily":             [dict(r) for r in daily],
        "top_workspaces":    [dict(r) for r in top_ws],
    })


# ── File History ──────────────────────────────────────────────────────────────
@app.route('/file_history')
@login_required
def file_history():
    db = get_db()
    rows = db.execute(
        "SELECT filename, filetype, chat_id, uploaded_at FROM uploaded_files WHERE username=? ORDER BY uploaded_at DESC LIMIT 50",
        (session['user'],)).fetchall()
    return jsonify({"files": [dict(r) for r in rows]})



# ── History Page ──────────────────────────────────────────────────────────────
@app.route('/history')
@login_required
def history_page():
    return render_template('history.html', username=session['user'])


@app.route('/api/history')
@login_required
def api_history():
    """Return paginated, filterable chat history for the history page."""
    page        = int(request.args.get('page', 1))
    per_page    = int(request.args.get('per_page', 20))
    search      = request.args.get('search', '').strip()
    pipeline    = request.args.get('pipeline', '').strip()
    workspace   = request.args.get('workspace', '').strip()
    date_from   = request.args.get('date_from', '').strip()
    date_to     = request.args.get('date_to', '').strip()
    offset      = (page - 1) * per_page

    db     = get_db()
    params = [session['user']]
    where  = ["username=?"]

    if search:
        where.append("(message LIKE ? OR response LIKE ?)")
        params += [f'%{search}%', f'%{search}%']
    if pipeline:
        where.append("routing LIKE ?")
        params.append(f'%{pipeline}%')
    if workspace:
        where.append("chat_id=?")
        params.append(workspace)
    if date_from:
        where.append("DATE(timestamp) >= ?")
        params.append(date_from)
    if date_to:
        where.append("DATE(timestamp) <= ?")
        params.append(date_to)

    where_clause = " AND ".join(where)

    total = db.execute(
        f"SELECT COUNT(*) as c FROM chats WHERE {where_clause}", params
    ).fetchone()['c']

    rows = db.execute(
        f"""SELECT id, chat_id, message, response, routing, response_time, pinned, timestamp
            FROM chats WHERE {where_clause}
            ORDER BY id DESC LIMIT ? OFFSET ?""",
        params + [per_page, offset]
    ).fetchall()

    # All distinct workspaces for filter dropdown
    workspaces = db.execute(
        "SELECT DISTINCT chat_id FROM chats WHERE username=? ORDER BY chat_id",
        (session['user'],)
    ).fetchall()

    return jsonify({
        "total":      total,
        "page":       page,
        "per_page":   per_page,
        "pages":      (total + per_page - 1) // per_page,
        "rows":       [dict(r) for r in rows],
        "workspaces": [r['chat_id'] for r in workspaces],
    })


# ── Analytics Page ────────────────────────────────────────────────────────────
@app.route('/analytics')
@login_required
def analytics_page():
    return render_template('analytics.html', username=session['user'])

# ── Static ────────────────────────────────────────────────────────────────────
@app.route('/style.css')
def style():
    return send_from_directory('static', 'style.css')

@app.route('/script.js')
def script_js():
    return send_from_directory('static', 'script.js')

if __name__ == '__main__':
    app.run(debug=True)